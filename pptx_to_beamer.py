import os
import argparse
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def escape_latex(text):
    """
    Escape special LaTeX characters in a string.
    """
    if not isinstance(text, str):
        return text
    # A more robust way to handle special characters is to replace them
    # with their text equivalents, but for this simple script,
    # we'll just remove characters that are known to cause issues.
    return text.replace('&', 'and') \
               .replace('%', '') \
               .replace('$', '') \
               .replace('#', '') \
               .replace('_', ' ') \
               .replace('{', '') \
               .replace('}', '') \
               .replace('~', '') \
               .replace('^', '') \
               .replace('\\', '') \
               .replace('\v', ' ') # Replace vertical tab with a space

def sanitize_filename(text):
    """
    Sanitize a string to be used as a filename.
    """
    if not isinstance(text, str):
        text = "untitled"
    text = text.lower()
    text = re.sub(r'[^a-z0-9_]+', '_', text)
    return text.strip('_')


def extract_pptx_content(pptx_path, output_dir):
    """
    Extracts text and images from a PPTX file.
    """
    prs = Presentation(pptx_path)
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)

    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_data = {'text': [], 'images': []}
        
        # Extract title
        slide_title = "untitled"
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
            slide_data['title'] = escape_latex(slide_title)

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Avoid duplicating title
                if shape.has_text_frame and shape.text != slide_data.get('title', ''):
                    slide_data['text'].append(escape_latex(shape.text))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                sanitized_title = sanitize_filename(slide_title)
                image_filename = f'slide_{i+1}_{sanitized_title}_image_{len(slide_data["images"]) + 1}.{image.ext}'
                image_path = os.path.join(images_dir, image_filename)
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                slide_data['images'].append(os.path.join('images', image_filename))

        slides_content.append(slide_data)

    # Attempt to get presentation metadata
    title = escape_latex(prs.core_properties.title)
    author = escape_latex(prs.core_properties.author)
    
    return slides_content, title, author

def generate_beamer_tex(slides_content, title, author, output_dir):
    """
    Generates a Beamer .tex file from the extracted content.
    """
    tex_file_path = os.path.join(output_dir, 'presentation.tex')

    with open(tex_file_path, 'w', encoding='utf-8') as f:
        f.write('\\documentclass{beamer}\n')
        f.write('\\usepackage{graphicx}\n')
        f.write('\\usepackage[utf8]{inputenc}\n')
        f.write('\\usepackage{ragged2e}\n\n')
        
        if title:
            f.write(f'\\title{{{title}}}\n')
        if author:
            f.write(f'\\author{{{author}}}\n')
        f.write('\\date{\\today}\n\n')
        
        f.write('\\begin{document}\n\n')
        
        if title:
            f.write('\\frame{\\titlepage}\n\n')

        for i, slide in enumerate(slides_content):
            f.write('\\begin{frame}[fragile]\n') # Use 'fragile' for verbatim
            
            if slide.get('title'):
                f.write(f'  \\frametitle{{{slide["title"]}}}\n')
            
            f.write('  \\begin{verbatim}\n')
            for text in slide['text']:
                f.write(f'    {text}\n\n')
            f.write('  \\end{verbatim}\n')

            for image_path in slide['images']:
                # Use a relative path for the image
                f.write(f'  \\includegraphics[width=0.8\\textwidth, height=0.6\\textheight, keepaspectratio]{{{image_path}}}\n')
            
            f.write('\\end{frame}\n\n')

        f.write('\\end{document}\n')

    return tex_file_path

def main():
    parser = argparse.ArgumentParser(description='Convert a PowerPoint presentation to a LaTeX Beamer presentation.')
    parser.add_argument('input_pptx', help='The input PPTX file.')
    parser.add_argument('output_dir', help='The directory to save the .tex file and images.')
    args = parser.parse_args()

    if not os.path.exists(args.input_pptx):
        print(f"Error: Input file not found at {args.input_pptx}")
        return

    os.makedirs(args.output_dir, exist_ok=True)

    slides_content, title, author = extract_pptx_content(args.input_pptx, args.output_dir)
    tex_file_path = generate_beamer_tex(slides_content, title, author, args.output_dir)

    print(f"Successfully generated {tex_file_path}")

if __name__ == '__main__':
    main()
